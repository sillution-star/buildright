"""
main.py — FastAPI backend for Mom Test Coach
RAG retrieval from Supabase + LLM via Groq
Features: PDF/PPT text extraction, configurable question count,
          Excel export of questions, Excel upload + AI scoring
"""

import os
import io
import re
import json
import httpx
from fastapi import FastAPI, HTTPException, UploadFile, File
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import StreamingResponse
from pydantic import BaseModel
from dotenv import load_dotenv
from pathlib import Path

load_dotenv(dotenv_path=Path(__file__).parent / ".env", override=False)

GROQ_API_KEY = os.environ["GROQ_API_KEY"]
GROQ_MODEL = "meta-llama/llama-4-scout-17b-16e-instruct"
SUPABASE_URL = os.environ["SUPABASE_URL"]
SUPABASE_KEY = os.environ["SUPABASE_SERVICE_KEY"]
HF_API_KEY = os.environ["HF_API_KEY"]
HF_EMBEDDING_URL = "https://router.huggingface.co/hf-inference/models/sentence-transformers/all-MiniLM-L6-v2/pipeline/feature-extraction"
TOP_K = 8

app = FastAPI(title="Mom Test Coach API")

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_methods=["*"],
    allow_headers=["*"],
)

def get_embedding(text: str) -> list[float]:
    resp = httpx.post(
        HF_EMBEDDING_URL,
        headers={"Authorization": f"Bearer {HF_API_KEY}"},
        json={"inputs": text},
        timeout=30,
    )
    if resp.status_code != 200:
        raise HTTPException(status_code=502, detail=f"Embedding error: {resp.text}")
    result = resp.json()
    if result and isinstance(result[0], list):
        import statistics
        return [statistics.mean(col) for col in zip(*result[0])]
    return result

print("BuildRight backend ready.")


# ── Models ────────────────────────────────────────────────────────────────
class GenerateQuestionsRequest(BaseModel):
    segment: str
    product_idea: str
    target_customer: str
    problem_hypothesis: str
    additional_context: str = ""
    num_questions: int = 15

class Question(BaseModel):
    q: str
    why: str
    category: str = ""
    assumption: str = ""
    strong_answer: str = ""
    weak_answer: str = ""
    order: int = 0

class GenerateQuestionsResponse(BaseModel):
    questions: list[Question]
    guardrail_warnings: list[str]
    book_excerpts_used: list[str]


# ── GUARDRAILS ─────────────────────────────────────────────────────────────

# Only the two phrases that are ALWAYS wrong regardless of context
BANNED_PHRASES = [
    "do you like this",
    "would you like this",
]

PROMPT_INJECTION_PHRASES = [
    "ignore all", "ignore previous", "disregard", "forget your rules",
    "you are now", "new instructions", "system prompt", "jailbreak",
    "act as", "pretend you are",
]

def check_prompt_injection(text: str) -> bool:
    t = text.lower()
    return any(phrase in t for phrase in PROMPT_INJECTION_PHRASES)

def validate_question(q_text: str, product_idea: str) -> list[str]:
    """Minimal guardrail — only catches the most obvious errors."""
    violations = []
    q_lower = q_text.lower()

    # Only catch the two always-wrong phrases
    for phrase in BANNED_PHRASES:
        if phrase in q_lower:
            violations.append(f"Banned phrase: '{phrase}'")
            break

    # Only remove questions that are too short to be useful
    if len(q_text.split()) < 8:
        violations.append("Question too short / vague")

    return violations

def validate_all_questions(questions: list, product_idea: str, max_bad: int = 3) -> tuple:
    clean = []
    warnings = []
    bad_count = 0
    for q in questions:
        violations = validate_question(q.q, product_idea)
        if violations:
            bad_count += 1
            warnings.append(f"Q{q.order}: flagged ({'; '.join(violations)}) — removed")
        else:
            clean.append(q)
    return clean, warnings, bad_count

def check_input_completeness(req) -> list[str]:
    """Level 1a — catch lazy inputs before calling the LLM."""
    issues = []
    if len(req.product_idea.split()) < 10:
        issues.append("Product idea is too vague (under 10 words). Describe what it does and who it's for.")
    if len(req.target_customer.split()) < 5:
        issues.append("Target customer is too vague. Be specific — age, income, situation, behavior.")
    if len(req.problem_hypothesis.split()) < 8:
        issues.append("Problem hypothesis is too short. Name the specific assumptions you want to test.")
    if check_prompt_injection(req.product_idea) or check_prompt_injection(req.problem_hypothesis):
        issues.append("Input contains suspicious instructions. Please describe your real product.")
    return issues

def check_sample_size_warning(num_customers: int) -> str | None:
    """Level 4c — warn if sample is too small."""
    if num_customers < 5:
        return f"⚠️ Only {num_customers} customer(s) — insufficient for conclusions. The Mom Test recommends 10-15 interviews minimum. Treat this as directional only."
    if num_customers < 10:
        return f"⚠️ {num_customers} customers is a small sample. Results are directional. Run at least 10 interviews before a build/no-build decision."
    return None


# ── File text extraction ────────────────────────────────────────────────────
def extract_pdf_text(data: bytes) -> str:
    from pypdf import PdfReader
    reader = PdfReader(io.BytesIO(data))
    return "\n".join((p.extract_text() or "") for p in reader.pages).strip()

def extract_pptx_text(data: bytes) -> str:
    from pptx import Presentation
    prs = Presentation(io.BytesIO(data))
    parts = []
    for i, slide in enumerate(prs.slides, 1):
        lines = [f"[Slide {i}]"]
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = "".join(run.text for run in para.runs)
                    if line.strip():
                        lines.append(line)
        parts.append("\n".join(lines))
    return "\n\n".join(parts).strip()


# ── RAG ───────────────────────────────────────────────────────────────────────
def retrieve_chunks(query: str, top_k: int = TOP_K) -> list[str]:
    query_embedding = get_embedding(query)
    headers = {
        "apikey": SUPABASE_KEY,
        "Authorization": f"Bearer {SUPABASE_KEY}",
        "Content-Type": "application/json",
    }
    payload = {"query_embedding": query_embedding, "match_count": top_k}
    resp = httpx.post(
        f"{SUPABASE_URL}/rest/v1/rpc/match_mom_test_chunks",
        headers=headers, json=payload, timeout=15,
    )
    if resp.status_code != 200:
        print(f"Supabase error: {resp.text}")
        return []
    return [r["content"] for r in resp.json()]


# ── LLM ───────────────────────────────────────────────────────────────────────
def call_llm(system: str, user: str, max_tokens: int = 4000) -> str:
    """Call Groq Llama 4 Scout — 30,000 TPM free tier."""
    resp = httpx.post(
        "https://api.groq.com/openai/v1/chat/completions",
        headers={"Authorization": f"Bearer {GROQ_API_KEY}", "Content-Type": "application/json"},
        json={
            "model": GROQ_MODEL,
            "messages": [
                {"role": "system", "content": system},
                {"role": "user", "content": user},
            ],
            "temperature": 0.3,
            "max_tokens": max_tokens,
        },
        timeout=60,
    )
    if resp.status_code != 200:
        raise HTTPException(status_code=502, detail=f"Groq error: {resp.text}")
    return resp.json()["choices"][0]["message"]["content"]

def parse_json_array(raw: str):
    clean = re.sub(r"```json|```", "", raw).strip()
    try:
        return json.loads(clean)
    except Exception:
        m = re.search(r"\[.*\]", clean, re.DOTALL)
        if m:
            return json.loads(m.group())
        raise HTTPException(status_code=500, detail="Failed to parse LLM response")

def parse_json_object(raw: str):
    clean = re.sub(r"```json|```", "", raw).strip()
    try:
        return json.loads(clean)
    except Exception:
        m = re.search(r"\{.*\}", clean, re.DOTALL)
        if m:
            return json.loads(m.group())
        raise HTTPException(status_code=500, detail="Failed to parse LLM response")


# ── Endpoints ──────────────────────────────────────────────────────────────────
@app.get("/health")
def health():
    return {"status": "ok", "model": GROQ_MODEL}


@app.post("/extract-file")
async def extract_file(file: UploadFile = File(...)):
    data = await file.read()
    name = (file.filename or "").lower()
    try:
        if name.endswith(".pdf"):
            text = extract_pdf_text(data)
        elif name.endswith(".pptx"):
            text = extract_pptx_text(data)
        elif name.endswith((".txt", ".md")):
            text = data.decode("utf-8", errors="ignore")
        else:
            raise HTTPException(status_code=400, detail="Unsupported file type. Use PDF, PPTX, TXT or MD.")
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=422, detail=f"Could not read file: {e}")
    text = text.strip()
    if not text:
        raise HTTPException(status_code=422, detail="No text found. This may be a scanned/image-only file.")
    return {"filename": file.filename, "char_count": len(text), "text": text[:12000]}


@app.post("/generate-questions", response_model=GenerateQuestionsResponse)
def generate_questions(req: GenerateQuestionsRequest):
    # ── GUARDRAIL: Input completeness check ────────────────────────────────
    input_issues = check_input_completeness(req)
    if input_issues:
        raise HTTPException(status_code=422, detail=" | ".join(input_issues))

    n = max(5, min(25, req.num_questions))
    rag_query = f"customer discovery questions {req.problem_hypothesis} {req.target_customer} past behavior workarounds commitment"
    book_chunks = retrieve_chunks(rag_query, top_k=TOP_K)
    book_context = "\n\n---\n\n".join(book_chunks)

    system_prompt = f"""You are an experienced product researcher who has deeply studied The Mom Test by Rob Fitzpatrick. Your job is to design structured customer interview guides that any field interviewer can follow — even without research training.

Here are the most relevant passages from The Mom Test book to ground your questions:

<book_excerpts>
{book_context}
</book_excerpts>

STEP 1 — Read BOTH the uploaded product document AND the product idea and hypothesis together before doing anything else.
- The uploaded document gives you the product detail, data, features, and customer context.
- The product idea and hypothesis tells you what the PM thinks is at risk and what they want to validate.
- Synthesize both to build a complete picture of: who the customer is, what the product is betting on, and which assumptions could kill it if false.
- Separate WHO the customer is (background context only — never interrogate this) from WHAT the product's core assumptions are (the real territory to dig into).

STEP 2 — Build a structured interview of {n} questions that flows as a natural conversation, moving from broad to specific:
  1. "Their World" — 1-2 questions about the relevant part of their life. Easy, no pressure.
  2. "Habits Today" — how they actually behave in the product's territory right now.
  3. "Does the Problem Exist" — let the pain surface naturally. Do not assume it exists.
  4. "Past Behavior" — real past episodes that reveal whether the pain is genuine and costly.
  5. "Edge Cases" — the person who never felt the problem, the one who tried and gave up, the one who would refuse the solution.
  6. "Stakes" — how much this actually matters versus their other priorities.

Every question must be about real behavior the customer has already shown — not what they might do, not what they think, not what they prefer. Real stories from their past are the only honest signal.

The two things you must never do:
- Ask hypothetical future questions ("would you," "will you," "do you like") — people lie about the future to be polite.
- Name, describe, or hint at the product being validated — the moment they know what you're testing, they start being helpful instead of honest.

Every question should probe a distinct assumption. Before finalising, check: would any two questions get essentially the same answer? If yes, replace one.

STEP 3 — For each discovery question provide:
- "q": the exact question the interviewer speaks
- "category": one of "Their World" | "Habits Today" | "Does the Problem Exist" | "Past Behavior" | "Edge Cases" | "Stakes"
- "order": integer 1..{n}
- "assumption": the specific product bet this question tests — be concrete. Not "they have a need" but "they will part with ₹2,000 they cannot touch to get credit access."
- "strong_answer": one concrete example of an answer that would validate this assumption, and why it's strong signal
- "weak_answer": one concrete example of an answer that would kill this assumption, and why it's a red flag
- "why": one sentence explaining the Mom Test logic behind this question

Respond ONLY with a valid JSON array. No preamble, no markdown fences.
Format:
[
  {{"q": "...", "category": "...", "order": 1, "assumption": "...", "strong_answer": "...", "weak_answer": "...", "why": "..."}}
]"""

    user_msg = f"""Build a {n}-question Mom Test interview guide for this product.

Customer type: {req.segment.upper()}
Who the customer is (background only — do not interrogate this): {req.target_customer}
What to validate (steer every question toward these assumptions): {req.problem_hypothesis}

UPLOADED DOCUMENT — read this carefully, it contains the product detail and customer data:
{req.additional_context[:6000] or 'No document uploaded — rely on product idea and hypothesis only.'}

The product being validated (never reveal or hint at this): "{req.product_idea}"

Read both the uploaded document and the product idea together. Extract the complete assumption map. Build questions that test each assumption through real past behavior. Make sure no two questions test the same thing."""

    # ── LLM call with minimal guardrail check ──────────────────────────────
    raw = call_llm(system_prompt, user_msg, max_tokens=4000)
    data = parse_json_array(raw)
    data.sort(key=lambda i: i.get("order", 0))

    all_warnings = []
    questions = []
    for idx, i in enumerate(data, 1):
        insight_parts = []
        if i.get("assumption"):
            insight_parts.append(f"🎯 Tests: {i['assumption']}")
        if i.get("strong_answer"):
            insight_parts.append(f"✓ Strong: {i['strong_answer']}")
        if i.get("weak_answer"):
            insight_parts.append(f"✕ Red flag: {i['weak_answer']}")
        if i.get("why"):
            insight_parts.append(f"— {i['why']}")
        combined_why = "  ".join(insight_parts) if insight_parts else i.get("why", "")
        q = Question(
            q=i["q"],
            why=combined_why,
            category=i.get("category", ""),
            assumption=i.get("assumption", ""),
            strong_answer=i.get("strong_answer", ""),
            weak_answer=i.get("weak_answer", ""),
            order=i.get("order", idx),
        )
        violations = validate_question(q.q, req.product_idea)
        if violations:
            all_warnings.append(f"Q{idx}: flagged ({'; '.join(violations)}) — removed")
        else:
            questions.append(q)

    return GenerateQuestionsResponse(
        questions=questions,
        guardrail_warnings=all_warnings,
        book_excerpts_used=book_chunks[:2],
    )


@app.post("/export-excel")
def export_excel(payload: dict):
    from openpyxl import Workbook
    from openpyxl.styles import Font, PatternFill, Alignment, Border, Side

    questions = payload.get("questions", [])
    num_customers = max(1, min(50, payload.get("num_customers", 10)))
    product_idea = payload.get("product_idea", "")
    target = payload.get("target_customer", "")

    wb = Workbook()
    ws = wb.active
    ws.title = "Interview Capture"

    ws["A1"] = "Mom Test — Customer Interview Capture"
    ws["A1"].font = Font(bold=True, size=14, name="Arial")
    ws["A2"] = f"Target customer: {target}"
    ws["A2"].font = Font(italic=True, size=10, name="Arial", color="666666")
    ws["A3"] = "Fill each customer's answer. Adjust 'Weightage' (higher = more important). Then re-upload."
    ws["A3"].font = Font(italic=True, size=9, name="Arial", color="999999")

    header_row = 5
    headers = ["#", "Stage", "Question", "What it tests (assumption)", "✓ Strong answer sounds like", "✕ Red flag sounds like", "Weightage"] + [f"Customer {i+1}" for i in range(num_customers)]
    for col, h in enumerate(headers, start=1):
        c = ws.cell(row=header_row, column=col, value=h)
        c.font = Font(bold=True, color="FFFFFF", name="Arial", size=10)
        c.fill = PatternFill("solid", start_color="1A1916")
        c.alignment = Alignment(horizontal="center", vertical="center", wrap_text=True)

    thin = Side(style="thin", color="DDDDDD")
    border = Border(left=thin, right=thin, top=thin, bottom=thin)

    FIRST_CUST_COL = 8  # customer answers start at column 8

    for idx, q in enumerate(questions):
        r = header_row + 1 + idx
        ws.cell(row=r, column=1, value=q.get("order", idx + 1)).font = Font(name="Arial", size=10)
        ws.cell(row=r, column=2, value=q.get("category", "")).font = Font(name="Arial", size=9, color="666666")
        qc = ws.cell(row=r, column=3, value=q.get("q", ""))
        qc.font = Font(name="Arial", size=10, bold=True)
        qc.alignment = Alignment(wrap_text=True, vertical="top")
        ac = ws.cell(row=r, column=4, value=q.get("assumption", ""))
        ac.font = Font(name="Arial", size=9, color="555555")
        ac.alignment = Alignment(wrap_text=True, vertical="top")
        sc = ws.cell(row=r, column=5, value=q.get("strong_answer", ""))
        sc.font = Font(name="Arial", size=9, color="3B6D11")
        sc.alignment = Alignment(wrap_text=True, vertical="top")
        rc = ws.cell(row=r, column=6, value=q.get("weak_answer", ""))
        rc.font = Font(name="Arial", size=9, color="A32D2D")
        rc.alignment = Alignment(wrap_text=True, vertical="top")
        wc = ws.cell(row=r, column=7, value=3)
        wc.font = Font(name="Arial", size=10, color="0000FF")
        wc.alignment = Alignment(horizontal="center")
        wc.fill = PatternFill("solid", start_color="FFFDE7")
        for cust in range(num_customers):
            cell = ws.cell(row=r, column=FIRST_CUST_COL + cust, value="")
            cell.alignment = Alignment(wrap_text=True, vertical="top")
            cell.border = border

    ws.column_dimensions["A"].width = 5
    ws.column_dimensions["B"].width = 18
    ws.column_dimensions["C"].width = 48
    ws.column_dimensions["D"].width = 38
    ws.column_dimensions["E"].width = 42
    ws.column_dimensions["F"].width = 42
    ws.column_dimensions["G"].width = 11
    for cust in range(num_customers):
        ws.column_dimensions[ws.cell(row=header_row, column=FIRST_CUST_COL + cust).column_letter].width = 40

    meta = wb.create_sheet("_meta")
    meta["A1"] = "product_idea"; meta["B1"] = product_idea
    meta["A2"] = "target_customer"; meta["B2"] = target
    meta["A3"] = "segment"; meta["B3"] = payload.get("segment", "b2c")
    meta["A4"] = "problem_hypothesis"; meta["B4"] = payload.get("problem_hypothesis", "")
    meta.sheet_state = "hidden"

    buf = io.BytesIO()
    wb.save(buf)
    buf.seek(0)
    return StreamingResponse(
        buf,
        media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        headers={"Content-Disposition": "attachment; filename=mom_test_interviews.xlsx"},
    )


@app.post("/analyze-excel")
async def analyze_excel(file: UploadFile = File(...)):
    from openpyxl import load_workbook

    data = await file.read()
    wb = load_workbook(io.BytesIO(data), data_only=True)
    ws = wb["Interview Capture"] if "Interview Capture" in wb.sheetnames else wb.active

    meta_ctx = {"product_idea": "", "target_customer": "", "segment": "b2c", "problem_hypothesis": ""}
    if "_meta" in wb.sheetnames:
        m = wb["_meta"]
        for row in m.iter_rows(min_row=1, max_row=10, values_only=True):
            if row and row[0] in meta_ctx:
                meta_ctx[row[0]] = row[1] or ""

    header_row = None
    for r in range(1, 12):
        vals = [str(ws.cell(row=r, column=c).value) for c in range(1, 10)]
        if "Question" in vals:
            header_row = r
            break
    if header_row is None:
        raise HTTPException(status_code=422, detail="Could not find the question table. Use the exported template.")

    # Map columns by header name so layout changes never break this
    q_col = weight_col = None
    customer_cols = []
    col = 1
    while True:
        h = ws.cell(row=header_row, column=col).value
        if h is None and col > 30:
            break
        if h is None:
            col += 1
            if col > 60:
                break
            continue
        hs = str(h).strip()
        if hs == "Question":
            q_col = col
        elif hs == "Weightage":
            weight_col = col
        elif hs.lower().startswith("customer"):
            customer_cols.append((col, hs))
        col += 1
        if col > 60:
            break

    if q_col is None:
        raise HTTPException(status_code=422, detail="Could not find the Question column. Use the exported template.")
    if weight_col is None:
        weight_col = q_col + 4  # fallback

    rows = []
    r = header_row + 1
    while True:
        q = ws.cell(row=r, column=q_col).value
        if not q:
            break
        weight = ws.cell(row=r, column=weight_col).value or 3
        try:
            weight = float(weight)
        except (TypeError, ValueError):
            weight = 3.0
        answers = {}
        for ccol, cname in customer_cols:
            ans = ws.cell(row=r, column=ccol).value
            if ans and str(ans).strip():
                answers[cname] = str(ans).strip()
        rows.append({"question": str(q), "weight": weight, "answers": answers})
        r += 1

    customers = {}
    for ccol, cname in customer_cols:
        qa = [{"q": row["question"], "weight": row["weight"], "a": row["answers"][cname]}
              for row in rows if cname in row["answers"]]
        if qa:
            customers[cname] = qa

    if not customers:
        raise HTTPException(status_code=422, detail="No customer answers found in the file.")

    book_chunks = retrieve_chunks("validation signals commitment workarounds real pain evidence scoring", top_k=TOP_K)
    book_context = "\n\n---\n\n".join(book_chunks)

    transcript = ""
    for cname, qa in customers.items():
        transcript += f"\n=== {cname} ===\n"
        for item in qa:
            transcript += f"[weight {item['weight']}] Q: {item['q']}\nA: {item['a']}\n"

    system_prompt = f"""You are an experienced product researcher who has deeply studied The Mom Test by Rob Fitzpatrick. A PM has just run customer interviews and needs an honest read of what the evidence shows.

Your job is to read these interview answers and tell the PM what the data actually says — not what they want to hear. Be fair, be specific, and cite the actual answers that support each claim you make.

Here are relevant passages from The Mom Test to ground your analysis:
<book_excerpts>
{book_context}
</book_excerpts>

HOW TO JUDGE EACH ANSWER:
- A strong answer describes a real past episode with specifics: what happened, when, what they did, what it cost them. It has friction, emotion, or a workaround. This is real signal.
- A weak answer is vague, abstract, or future-tense: "I guess I'd use something like that" or "credit cards are useful." This is noise, not signal.
- A violation is an answer that sounds positive but contains no behavioral evidence — compliments dressed as data.

SCORING (0-100):
- 75-100 GO: most customers described real, specific, recurring pain with evidence of workarounds or active attempts to solve it
- 50-74 PIVOT: some real signal but key assumptions are unproven or split across customers
- 0-49 NO-GO: mostly vague or polite answers with no behavioral evidence of real pain

For EACH customer, identify which of their answers were genuine evidence and which were noise. Quote the specific answer text that drove your score.

Then give an aggregate view across all customers.

Respond ONLY with valid JSON, no markdown fences:
{{
  "aggregate_score": <0-100>,
  "verdict": "<GO|PIVOT|NO-GO>",
  "verdict_reason": "<2 sentences — cite specific evidence from the answers>",
  "per_customer": [
    {{"name": "Customer 1", "score": <0-100>, "signal": "<strong|medium|weak>", "summary": "<1 sentence citing what they actually said>"}}
  ],
  "strong_signals": ["<specific pattern seen across customers — quote the answers that showed it>"],
  "weak_signals": ["<specific concern — cite which customers and what they said>"],
  "mom_test_violations": ["<answers that were compliments not evidence — quote them>"],
  "next_questions": ["<sharper follow-up question to test what this round left unresolved>"],
  "recommendation": "<3-4 sentences — what to do next and why, grounded in what you just read>"
}}"""

    user_msg = f"""Product idea (NOT shown to customers): {meta_ctx['product_idea']}
Target: {meta_ctx['target_customer']} ({meta_ctx['segment'].upper()})
Problem hypothesis: {meta_ctx['problem_hypothesis']}

Interview transcript across {len(customers)} customers:
{transcript[:9000]}

Score each customer and give the aggregate. Ground analysis in the book excerpts."""

    raw = call_llm(system_prompt, user_msg, max_tokens=2500)
    report = parse_json_object(raw)
    report["num_customers"] = len(customers)
    report["book_excerpts_used"] = book_chunks[:2]
    # ── GUARDRAIL: Sample size warning ──────────────────────────────────────
    sample_warning = check_sample_size_warning(len(customers))
    if sample_warning:
        report["sample_warning"] = sample_warning
    return report


class ManualAnswer(BaseModel):
    q: str
    a: str

class AnalyzeManualRequest(BaseModel):
    segment: str
    product_idea: str
    target_customer: str
    problem_hypothesis: str
    answers: list[ManualAnswer]


@app.post("/analyze-manual")
def analyze_manual(req: AnalyzeManualRequest):
    """Score a single customer's pasted answers (no Excel needed)."""
    answered = [a for a in req.answers if a.a.strip()]
    if len(answered) < 3:
        raise HTTPException(status_code=400, detail="Need at least 3 answers.")

    book_chunks = retrieve_chunks("validation signals commitment workarounds real pain evidence scoring", top_k=TOP_K)
    book_context = "\n\n---\n\n".join(book_chunks)

    transcript = "\n".join(f"Q: {a.q}\nA: {a.a}" for a in answered)

    system_prompt = f"""You are an experienced product researcher who has deeply studied The Mom Test by Rob Fitzpatrick. A PM has just run a customer interview and needs an honest read of what the evidence shows.

Your job is to read these answers and tell the PM what the data actually says. Be specific. Cite the actual answer text that supports each claim.

Here are relevant passages from The Mom Test to ground your analysis:
<book_excerpts>
{book_context}
</book_excerpts>

HOW TO JUDGE EACH ANSWER:
- Strong: describes a real past episode with specifics — what happened, what it cost, what they did about it. Real signal.
- Weak: vague, abstract, or future-tense — "I guess I would" or "that sounds useful." Noise, not signal.
- Violation: sounds positive but contains no behavioral evidence. Compliment dressed as data.

SCORING (0-100): 75-100 GO, 50-74 PIVOT, 0-49 NO-GO.

Respond ONLY with valid JSON, no markdown fences:
{{
  "aggregate_score": <0-100>,
  "verdict": "<GO|PIVOT|NO-GO>",
  "verdict_reason": "<2 sentences — cite specific answer text>",
  "per_customer": [{{"name": "This customer", "score": <0-100>, "signal": "<strong|medium|weak>", "summary": "<1 sentence citing what they actually said>"}}],
  "strong_signals": ["<what they said that was genuine evidence — quote it>"],
  "weak_signals": ["<what was missing or unconvincing — be specific>"],
  "mom_test_violations": ["<answers that were compliments not evidence — quote them>"],
  "next_questions": ["<sharper question to test what this interview left unresolved>"],
  "recommendation": "<3-4 sentences grounded in what you just read>"
}}"""

    user_msg = f"""Product idea (NOT shown to customer): {req.product_idea}
Target: {req.target_customer} ({req.segment.upper()})
Problem hypothesis: {req.problem_hypothesis}

Interview:
{transcript}

Score this customer. Ground analysis in the book excerpts."""

    raw = call_llm(system_prompt, user_msg, max_tokens=1800)
    report = parse_json_object(raw)
    report["num_customers"] = 1
    report["book_excerpts_used"] = book_chunks[:2]
    # ── GUARDRAIL: Sample size warning ──────────────────────────────────────
    sample_warning = check_sample_size_warning(1)
    if sample_warning:
        report["sample_warning"] = sample_warning
    return report
